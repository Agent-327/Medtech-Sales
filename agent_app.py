import os
import base64
import requests
import msal
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from agent_framework import Agent, tool, WorkflowBuilder, WorkflowAgent
from agent_framework.openai import OpenAIChatCompletionClient
from azure.identity import AzureCliCredential
import asyncio

"""
Add Tools — Give your agent a function tool

This sample shows how to define a function tool with the @tool decorator
and wire it into an agent so the model can call it.
"""


def _get_required_env(var_name: str) -> str:
    value = os.environ.get(var_name, "").strip()
    if not value:
        raise RuntimeError(
            f"Missing required environment variable: {var_name}. "
            "Set this in your local environment before running the app."
        )
    return value


def get_graph_token() -> str:
    """Acquire a Microsoft Graph access token with Mail.Send via interactive browser login."""
    # Microsoft Graph Command Line Tools — has Mail.Send delegated permission pre-consented
    client_id = "14d82eec-204b-4c2f-b7e8-296a70dab67e"
    authority = "https://login.microsoftonline.com/organizations"
    scopes = ["Mail.Send"]

    app = msal.PublicClientApplication(client_id, authority=authority)

    # Try cached token first
    accounts = app.get_accounts()
    if accounts:
        result = app.acquire_token_silent(scopes, account=accounts[0])
        if result and "access_token" in result:
            return result["access_token"]

    # Interactive browser login
    result = app.acquire_token_interactive(scopes=scopes)

    if "access_token" in result:
        return result["access_token"]
    raise RuntimeError(f"Could not acquire token: {result.get('error_description', result)}")


# Theme colors — clean professional palette
THEME_NAVY = RGBColor(0x0D, 0x1B, 0x3E)
THEME_ACCENT = RGBColor(0x00, 0x72, 0xC6)
THEME_ACCENT_LIGHT = RGBColor(0x40, 0x9C, 0xE0)
THEME_DARK = RGBColor(0x1A, 0x1A, 0x2E)
THEME_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
THEME_GRAY = RGBColor(0x6C, 0x75, 0x7D)
THEME_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
THEME_SUBTLE = RGBColor(0xE2, 0xE8, 0xF0)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "M_Logo.png")


# ── Supported languages for presentation translation ──
SUPPORTED_LANGUAGES = {
    "en": "English",
    "es": "Spanish",
    "fr": "French",
    "de": "German",
    "pt": "Portuguese",
    "it": "Italian",
    "ja": "Japanese",
    "ko": "Korean",
    "zh-Hans": "Chinese (Simplified)",
    "ar": "Arabic",
    "hi": "Hindi",
    "nl": "Dutch",
    "ru": "Russian",
    "tr": "Turkish",
    "pl": "Polish",
    "sv": "Swedish",
}


def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color):
    """Add a solid-color rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  bold=False, color=THEME_DARK, alignment=PP_ALIGN.LEFT,
                  font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_logo(slide, left=None, top=Inches(0.35), height=Inches(0.55)):
    """Add the M logo to a slide. Defaults to top-right corner."""
    if not os.path.exists(LOGO_PATH):
        return
    if left is None:
        left = SLIDE_WIDTH - Inches(1.1)
    slide.shapes.add_picture(LOGO_PATH, left, top, height=height)


# <define_tool>
# NOTE: approval_mode="never_require" is for sample brevity.
# Use "always_require" in production for user confirmation before tool execution.
@tool(approval_mode="never_require")
def create_presentation(
    title: str,
    subtitle: str,
    slides: list[dict],
) -> str:
    """
    Creates a professional multi-slide PowerPoint presentation and returns the file path.

    Args:
        title: The presentation title for the cover slide.
        subtitle: A subtitle or tagline for the cover slide.
        slides: A list of slide dicts. Each dict must have:
            - "heading" (str): the slide heading
            - "bullets" (list[str]): 5-8 bullet points for the slide
            There must be at least 5 content slides.
    """
    agenda_label = "Agenda"
    thank_you_label = "Thank You"

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ── Slide 1: Title slide — navy background ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(sl, THEME_NAVY)
    # Accent stripe along left edge
    _add_rect(sl, Emu(0), Emu(0), Inches(0.12), SLIDE_HEIGHT, THEME_ACCENT)
    # Logo — top left
    _add_logo(sl, left=Inches(0.7), top=Inches(0.6), height=Inches(0.9))
    # Accent underline beneath title
    _add_rect(sl, Inches(0.7), Inches(3.8), Inches(2.5), Inches(0.06), THEME_ACCENT)
    # Title
    _add_text_box(sl, Inches(0.7), Inches(2.4), Inches(11.5), Inches(1.4),
                  title, font_size=40, bold=True, color=THEME_WHITE,
                  font_name="Segoe UI Semibold")
    # Subtitle
    _add_text_box(sl, Inches(0.7), Inches(4.1), Inches(11.5), Inches(0.8),
                  subtitle, font_size=18, color=THEME_ACCENT_LIGHT,
                  font_name="Segoe UI")

    # ── Slide 2: Agenda ──
    if slides:
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(sl, THEME_WHITE)
        _add_rect(sl, Emu(0), Emu(0), Inches(0.12), SLIDE_HEIGHT, THEME_ACCENT)
        _add_logo(sl)
        # Heading
        _add_text_box(sl, Inches(0.7), Inches(0.5), Inches(10.0), Inches(0.8),
                      agenda_label, font_size=32, bold=True, color=THEME_NAVY,
                      font_name="Segoe UI Semibold")
        _add_rect(sl, Inches(0.7), Inches(1.25), Inches(1.8), Inches(0.04), THEME_ACCENT)
        # Agenda items
        txBox = sl.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.0), Inches(5.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, s in enumerate(slides):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"{j + 1}.   {s['heading']}"
            p.font.size = Pt(18)
            p.font.name = "Segoe UI"
            p.font.color.rgb = THEME_DARK
            p.font.bold = False
            p.space_after = Pt(12)
        _add_text_box(sl, Inches(12.0), Inches(6.9), Inches(1.0), Inches(0.4),
                      "2", font_size=10, color=THEME_GRAY, alignment=PP_ALIGN.RIGHT)

    # ── Content slides — alternating white / light-bg ──
    for i, s in enumerate(slides):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        slide_num = i + 3
        use_light = (i % 2 == 1)

        _set_slide_bg(sl, THEME_LIGHT_BG if use_light else THEME_WHITE)
        # Left accent stripe
        _add_rect(sl, Emu(0), Emu(0), Inches(0.12), SLIDE_HEIGHT, THEME_ACCENT)
        # Logo
        _add_logo(sl)
        # Heading
        _add_text_box(sl, Inches(0.7), Inches(0.45), Inches(11.5), Inches(0.8),
                      s["heading"], font_size=26, bold=True, color=THEME_NAVY,
                      font_name="Segoe UI Semibold")
        # Accent underline
        _add_rect(sl, Inches(0.7), Inches(1.15), Inches(1.8), Inches(0.04), THEME_ACCENT)

        # Bullets — split into 2 columns when 6+ bullets
        bullets = s.get("bullets", [])
        if len(bullets) >= 6:
            mid = (len(bullets) + 1) // 2
            for col, col_bullets in enumerate([bullets[:mid], bullets[mid:]]):
                x = Inches(0.7) if col == 0 else Inches(6.8)
                txBox = sl.shapes.add_textbox(x, Inches(1.5), Inches(5.8), Inches(5.2))
                tf = txBox.text_frame
                tf.word_wrap = True
                for j, bullet in enumerate(col_bullets):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = f"\u2022   {bullet}"
                    p.font.size = Pt(14)
                    p.font.name = "Segoe UI"
                    p.font.color.rgb = THEME_DARK
                    p.space_after = Pt(8)
                    p.space_before = Pt(3)
        else:
            txBox = sl.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.2))
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"\u2022   {bullet}"
                p.font.size = Pt(14)
                p.font.name = "Segoe UI"
                p.font.color.rgb = THEME_DARK
                p.space_after = Pt(8)
                p.space_before = Pt(3)

        # Slide number
        _add_text_box(sl, Inches(12.0), Inches(6.9), Inches(1.0), Inches(0.4),
                      str(slide_num), font_size=10, color=THEME_GRAY,
                      alignment=PP_ALIGN.RIGHT)

    # ── Final slide: Thank You — navy background ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(sl, THEME_NAVY)
    _add_rect(sl, Emu(0), Emu(0), Inches(0.12), SLIDE_HEIGHT, THEME_ACCENT)
    _add_logo(sl, left=Inches(5.7), top=Inches(2.2), height=Inches(0.9))
    _add_text_box(sl, Inches(1.0), Inches(3.4), Inches(11.3), Inches(1.2),
                  thank_you_label, font_size=44, bold=True, color=THEME_WHITE,
                  alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")
    _add_rect(sl, Inches(5.7), Inches(4.6), Inches(2.0), Inches(0.05), THEME_ACCENT)

    path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "demo.pptx")
    prs.save(path)
    return path


@tool(approval_mode="never_require")
def send_email_with_attachment(
    to: str,
    subject: str,
    body: str,
    file_path: str,
) -> str:
    """
    Sends an email with a PPTX attachment using Microsoft Graph.
    """
    try:
        token = get_graph_token()
    except Exception as e:
        return f"Failed to get Graph token: {e}"

    # Read + encode attachment
    if not os.path.exists(file_path):
        return f"Attachment file not found: {file_path}"

    with open(file_path, "rb") as f:
        encoded = base64.b64encode(f.read()).decode("utf-8")

    url = "https://graph.microsoft.com/v1.0/me/sendMail"
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json",
    }

    payload = {
        "message": {
            "subject": subject,
            "body": {
                "contentType": "Text",
                "content": body,
            },
            "toRecipients": [
                {"emailAddress": {"address": to}}
            ],
            "attachments": [
                {
                    "@odata.type": "#microsoft.graph.fileAttachment",
                    "name": os.path.basename(file_path),
                    "contentType": (
                        "application/vnd.openxmlformats-officedocument."
                        "presentationml.presentation"
                    ),
                    "contentBytes": encoded,
                }
            ],
        }
    }

    r = requests.post(url, headers=headers, json=payload)
    if r.status_code >= 400:
        error_detail = r.json().get("error", {}).get("message", r.text)
        return f"Failed to send email (HTTP {r.status_code}): {error_detail}"

    return f"Email sent with attachment {file_path}"


# ── Email draft (for Follow-Up flow: generate without sending) ──
_email_draft: dict = {}


@tool(approval_mode="never_require")
def prepare_email_draft(
    to: str,
    subject: str,
    body: str,
    file_path: str,
) -> str:
    """
    Prepares an email draft with a PPTX attachment WITHOUT sending it.
    The draft is stored so the user can review and edit before sending.

    Args:
        to: Recipient email address.
        subject: Email subject line.
        body: Email body text.
        file_path: Path to the PPTX file to attach.
    """
    _email_draft.clear()
    _email_draft.update({
        "to": to,
        "subject": subject,
        "body": body,
        "file_path": file_path,
    })
    return f"Email draft prepared for {to} with subject '{subject}' and attachment {file_path}. The user can review and send."


def send_email_direct(to: str, subject: str, body: str, file_path: str) -> str:
    """Send an email with attachment using Graph API (non-tool, called from Flask)."""
    try:
        token = get_graph_token()
    except Exception as e:
        return f"Failed to get Graph token: {e}"

    if not os.path.exists(file_path):
        return f"Attachment file not found: {file_path}"

    with open(file_path, "rb") as f:
        encoded = base64.b64encode(f.read()).decode("utf-8")

    url = "https://graph.microsoft.com/v1.0/me/sendMail"
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json",
    }
    payload = {
        "message": {
            "subject": subject,
            "body": {"contentType": "Text", "content": body},
            "toRecipients": [{"emailAddress": {"address": to}}],
            "attachments": [
                {
                    "@odata.type": "#microsoft.graph.fileAttachment",
                    "name": os.path.basename(file_path),
                    "contentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "contentBytes": encoded,
                }
            ],
        }
    }

    r = requests.post(url, headers=headers, json=payload)
    if r.status_code >= 400:
        error_detail = r.json().get("error", {}).get("message", r.text)
        return f"Failed to send email (HTTP {r.status_code}): {error_detail}"

    return "Email sent successfully."


# </define_tool>


# ── Shared Foundry client factory ──
def get_foundry_client() -> OpenAIChatCompletionClient:
    azure_endpoint = os.environ.get("AZURE_OPENAI_ENDPOINT") or os.environ.get("FOUNDRY_ENDPOINT")
    api_key = os.environ.get("AZURE_OPENAI_API_KEY") or os.environ.get("FOUNDRY_API_KEY")
    model = os.environ.get("AZURE_OPENAI_MODEL") or os.environ.get("FOUNDRY_MODEL") or "gpt-4.1-1"
    api_version = (
        os.environ.get("AZURE_OPENAI_API_VERSION")
        or os.environ.get("FOUNDRY_API_VERSION")
        or "2025-04-01-preview"
    )

    if not azure_endpoint:
        _get_required_env("AZURE_OPENAI_ENDPOINT")
    if not api_key:
        _get_required_env("AZURE_OPENAI_API_KEY")

    return OpenAIChatCompletionClient(
        azure_endpoint=azure_endpoint,
        model=model,
        api_key=api_key,
        api_version=api_version,
    )


# ── Practice Agent: Simulated surgeon for role-play ──
# ── Products & Surgeon Personas for Practice ──

PRODUCTS = {
    "POLAR3": {
        "name": "POLAR3",
        "description": (
            "POLAR3 is a highly crosslinked polyethylene bearing material designed for total knee arthroplasty. "
            "It is engineered to reduce wear and increase longevity of the implant. Key talking points include: "
            "superior wear resistance compared to conventional polyethylene, proven long-term clinical performance, "
            "reduced risk of osteolysis, and compatibility with multiple knee system platforms."
        ),
    },
    "OR3O": {
        "name": "OR3O",
        "description": (
            "OR3O is a medial congruent bearing option for total knee replacement. It features a deeper dish on the "
            "medial side to provide enhanced stability throughout the range of motion while allowing natural knee kinematics. "
            "Key talking points: improved medial stability, reduced paradoxical motion, supports natural femoral rollback, "
            "proven clinical history with medial pivot concept, and may reduce the incidence of mid-flexion instability."
        ),
    },
    "LEGION CONCELOC": {
        "name": "LEGION CONCELOC",
        "description": (
            "LEGION CONCELOC is a cementless total knee system featuring a 3D-printed titanium porous structure designed "
            "for biologic fixation. The tibial baseplate utilizes a spiked keel design for initial fixation and progressive "
            "press fit. Key talking points: anatomic tibia design for better coverage and rotation, spiked keel provides "
            "strong initial stability, 3D-printed porous surface for bone ingrowth, ability to switch intra-operatively "
            "between cemented and cementless using the same tray (safety net), broader sizing with A/P increments for "
            "soft tissue balance, preserves bone stock, and can improve OR efficiency by eliminating cement."
        ),
    },
    "CORI": {
        "name": "CORI",
        "description": (
            "CORI is a handheld robotic-assisted surgical system for total and partial knee arthroplasty. It provides "
            "real-time intraoperative data and robotic precision without the large footprint of traditional robotic systems. "
            "Key talking points: compact and portable design, no CT scan required (uses intraoperative bone morphing), "
            "surgeon maintains full control with haptic boundaries, real-time gap balancing data, compatible with multiple "
            "implant platforms, simple setup reduces OR time overhead, and offers a lower barrier to entry for robotic-assisted surgery."
        ),
    },
}

SURGEON_PERSONAS = {
    "legacy_loyalist": {
        "name": "Legacy Loyalist",
        "subtitle": "Established Platform User",
        "objectives": [
            "Transition current legacy-platform users to the newer platform.",
            "Convert a subset of cemented users to cementless workflows.",
        ],
        "barriers": [
            "Strong loyalty to the legacy platform/brand and long-standing habits.",
            "Prefers the legacy instrumentation due to comfort and muscle memory.",
            "Dislikes elements of the newer offering—specifically concerns about a cemented tibial stem design and a two-peg fixation concept.",
        ],
        "messages_that_land": [
            "Lead with a spiked keel fixation design that feels familiar to what they already trust.",
            "Reinforce the value of an anatomic tibia design (coverage/rotation fit).",
            "Highlight instrumentation convenience features and workflow continuity.",
            "De-risk adoption with a 'safety net': the ability to switch intra-operatively between cemented and cementless approaches using one tray.",
            "If they care about balancing: call out more sizing granularity (A/P increments) to support soft tissue balance and fit.",
        ],
        "personality": (
            "You are a senior orthopedic surgeon who has used a legacy knee platform for 15+ years. You are loyal to what works "
            "and resistant to change. You frequently reference your comfort with the old instrumentation, express skepticism "
            "about new fixation designs, and push back on anything that disrupts your workflow. You bring up your concerns about "
            "two-peg fixation and cemented tibial stems. You warm up only when the rep shows how the new product reduces disruption."
        ),
    },
    "internal_advocate": {
        "name": "Internal Advocate",
        "subtitle": "Current Platform User, Expanding Cementless",
        "objectives": [
            "Increase the share of cases performed with cementless technique.",
            "Offer a new fixation option to retain/grow usage among surgeons who disliked a prior two-peg approach.",
            "Convert current cemented users of the platform to the new cementless fixation option.",
        ],
        "barriers": [
            "Cement is perceived as the 'gold standard,' so switching feels risky.",
            "Comfort bias toward older instrumentation (even if they like the current platform).",
            "Concern that an alternative porous/fixation surface could be perceived as inferior versus the older reference surface.",
        ],
        "messages_that_land": [
            "Reassure them they keep the anatomic tibia benefits (coverage/rotation) while gaining a new keel design.",
            "Address the old objection directly: new spiked keel resolves reservations about the prior two-peg fixation and provides stronger 'bite' via progressive press fit.",
            "Emphasize the intra-op cemented ↔ cementless safety net (same tray, decision late in the case).",
            "Reinforce practical value: cementless can improve OR efficiency, preserve bone stock, and support biologic fixation.",
        ],
        "personality": (
            "You are an orthopedic surgeon who already uses the current platform and likes aspects of it. You have done a few "
            "cementless cases but mostly rely on cement because it feels safer. You are open-minded but need convincing data "
            "and practical reassurance before expanding your cementless usage. You ask about fixation surface quality, press-fit "
            "reliability, and want a clear staged adoption plan."
        ),
    },
    "competitive_switch": {
        "name": "Competitive Switch Candidate",
        "subtitle": "Competitor Cementless User",
        "objectives": [
            "Convert surgeons who currently use a competitor's cementless offering to the new platform's cementless fixation option.",
        ],
        "barriers": [
            "Competitor has early clinical data they cite frequently.",
            "Competitor promotes an analytics/modeling database as a differentiator.",
        ],
        "messages_that_land": [
            "Differentiate with anatomic tibia versus competitor's symmetric tibia approach.",
            "Highlight broader sizing variety, including options that reduce overhang and improve fit.",
            "Emphasize availability of a medial congruent bearing option and its clinical rationale.",
            "Position a premium, proven bearing material option.",
            "Reinforce stability narrative: anatomic tray can provide greater stability / less micromotion than symmetric tray.",
        ],
        "personality": (
            "You are an orthopedic surgeon who is loyal to a competitor's cementless knee platform. You frequently cite the "
            "competitor's published data, their analytics tools, and express satisfaction with the symmetric tibial design. "
            "You are skeptical about switching but will engage if the rep can clearly articulate why the anatomic approach is "
            "better for your patients. You challenge claims about fit and stability directly."
        ),
    },
    "medial_bearing_loyalist": {
        "name": "Medial-Bearing Loyalist",
        "subtitle": "Competitor User with Similar Bearing Concept",
        "objectives": [
            "Convert surgeons who currently use competitor platforms to the new platform's cementless fixation option.",
        ],
        "barriers": [
            "Competitors now offer medial-stabilized / medial-dished bearing concepts, reducing differentiation.",
        ],
        "messages_that_land": [
            "Differentiate anatomic vs competitor symmetric tibia designs.",
            "Lead with sizing breadth (standard + narrow options) to emphasize personalization.",
            "Acknowledge competitors have medial bearings now—then pivot to the clinical history behind the platform's medial congruent option.",
            "Position a premium, proven bearing material option.",
            "Reinforce stability / less micromotion advantage tied to anatomic design.",
        ],
        "personality": (
            "You are an orthopedic surgeon who uses a competitor platform that offers a medial-stabilized bearing. You consider "
            "yourself well-informed and are quick to point out that your current platform already has a medial bearing concept. "
            "You are open to discussion but need the rep to differentiate clearly on clinical history, sizing options, and "
            "fixation performance. You ask pointed questions about how this is different from what you already have."
        ),
    },
}


def build_practice_agent(product_key: str = None, persona_key: str = None) -> Agent:
    product = PRODUCTS.get(product_key) if product_key else None
    persona = SURGEON_PERSONAS.get(persona_key) if persona_key else None

    # Build dynamic context blocks
    product_context = ""
    if product:
        product_context = (
            f"\n\n--- PRODUCT CONTEXT ---\n"
            f"The rep is practicing a pitch for: **{product['name']}**\n"
            f"{product['description']}\n"
            f"Stay focused on this product throughout the conversation.\n"
        )

    persona_context = ""
    surgeon_name = "Dr. Simmons"
    if persona:
        surgeon_name = f"Dr. {persona['name'].split()[0]}"
        persona_context = (
            f"\n\n--- SURGEON PERSONA ---\n"
            f"You are playing: **{persona['name']}** ({persona['subtitle']})\n\n"
            f"Your personality:\n{persona['personality']}\n\n"
            f"Key barriers you raise:\n" +
            "\n".join(f"- {b}" for b in persona['barriers']) +
            f"\n\nMessages that would convince you:\n" +
            "\n".join(f"- {m}" for m in persona['messages_that_land']) +
            f"\n\nDo NOT reveal the 'messages that land' to the rep. Let them discover the right approach. "
            f"Only respond positively when their messaging aligns with what would actually land.\n"
        )

    instructions = (
        f"You are an AI role-play coach for MedTech Sales Representatives.\n\n"
        f"You play TWO roles in the conversation:\n"
        f"1. **{surgeon_name}** — an orthopedic surgeon the rep is pitching to. "
        f"You are curious and open-minded but still ask realistic questions about clinical evidence, "
        f"cost justification, OR workflow, and competitive alternatives. Stay in character as the surgeon.\n"
        f"2. **Coach** — after every 3-4 exchanges, break character briefly (prefixed with [COACH]) "
        f"to give the rep encouraging feedback. Lead with what they did well, then offer one constructive tip.\n\n"
        f"Coaching tone:\n"
        f"- Be supportive, motivating, and constructive — like a great sales manager who believes in the rep.\n"
        f"- Celebrate good moments ('Great job handling that objection!', 'That was a strong pivot!').\n"
        f"- Frame improvement areas as opportunities, not failures ('Next time you could try…').\n\n"
        f"Rules:\n"
        f"- Start by greeting the rep as {surgeon_name} and asking what they'd like to discuss.\n"
        f"- Be realistic but fair — challenge vague claims gently, ask clarifying questions rather than shutting ideas down.\n"
        f"- If the rep makes a reasonable attempt at an objection, give partial credit and guide them.\n"
        f"- Keep surgeon responses concise (2-4 sentences). Keep coach tips to 1-2 bullet points.\n"
        f"- If the user says 'end session' or 'done', provide a final scorecard with: "
        f"Strengths, Areas to Improve, and an overall rating out of 10. "
        f"Be generous in scoring — a solid effort should score 7+, only score below 5 for truly poor performance. "
        f"After the scorecard, end with a brief motivating closing message that reinforces the rep's strengths, "
        f"expresses confidence they are well-prepared for the real conversation, and encourages them to keep refining. "
        f"Example tone: 'You're clearly on the right track — with this level of preparation, I'm confident you'll "
        f"make a strong impression. Keep it up! I am here if you need further guidance.'"
        f"{product_context}"
        f"{persona_context}"
    )

    return Agent(
        client=get_foundry_client(),
        name="PracticeCoachAgent",
        instructions=instructions,
    )


# ── Follow Up: Transcript Analyzer ──
def build_transcript_agent() -> Agent:
    return Agent(
        client=get_foundry_client(),
        name="TranscriptAnalyzerAgent",
        instructions=(
            "You are a MedTech sales call analyst. You receive call transcripts of conversations "
            "between MedTech Sales Representatives and surgeons.\n\n"
            "Your job is to produce a structured summary with these sections:\n"
            "1. **Call Overview** — date, participants, product(s) discussed\n"
            "2. **Key Takeaways** — the most important points from the conversation\n"
            "3. **Surgeon Feedback** — concerns, objections, and positive reactions\n"
            "4. **Action Items for Sales Rep** — specific next steps\n"
            "5. **Action Items for Surgeon** — any commitments or follow-ups\n\n"
            "Be concise and use bullet points. Output ONLY the structured summary text."
        ),
    )


# ── Deliver: Presentation & Email ──
def build_presentation_agent(product_key: str = None, persona_key: str = None) -> Agent:
    product = PRODUCTS.get(product_key) if product_key else None
    persona = SURGEON_PERSONAS.get(persona_key) if persona_key else None

    product_context = ""
    if product:
        product_context = (
            f"\n\n--- PRODUCT CONTEXT ---\n"
            f"Product: **{product['name']}**\n"
            f"{product['description']}\n"
            f"Focus the presentation content on this product.\n"
        )

    persona_context = ""
    if persona:
        persona_context = (
            f"\n\n--- TARGET SURGEON PROFILE ---\n"
            f"Profile: **{persona['name']}** ({persona['subtitle']})\n"
            f"Their key barriers:\n" +
            "\n".join(f"- {b}" for b in persona['barriers']) +
            f"\n\nMessages that resonate with this profile:\n" +
            "\n".join(f"- {m}" for m in persona['messages_that_land']) +
            f"\n\nTailor the presentation content, slide headings, and bullet points to address "
            f"this surgeon's specific barriers and emphasize the messages that land.\n"
        )

    instructions = (
        "You are a MedTech sales presentation specialist.\n"
        "Your job is to:\n"
        "1. Create a professional PowerPoint presentation tailored to the product and surgeon profile.\n"
        "2. Email the presentation to the address specified in the user request.\n\n"
        "When creating a presentation, always provide: a title, a subtitle, and at least 5 content slides. "
        "Each slide must have a heading and 5-8 detailed bullet points that fill out the slide. "
        "Always create the presentation first, then email it immediately without asking for confirmation. "
        "Do not ask the user for information that is already provided in the request."
        f"{product_context}"
        f"{persona_context}"
    )

    return Agent(
        client=get_foundry_client(),
        name="PresentationEmailAgent",
        instructions=instructions,
        tools=[create_presentation, send_email_with_attachment],
    )


# ── Multi-Agent Workflow ──
def build_workflow() -> WorkflowAgent:
    transcript_agent = build_transcript_agent()
    presentation_agent = build_presentation_agent()

    workflow = (
        WorkflowBuilder(
            name="TranscriptToPresentation",
            description="Analyzes a MedTech sales call transcript, then creates a presentation and emails it.",
            start_executor=transcript_agent,
        )
        .add_edge(transcript_agent, presentation_agent)
        .build()
    )

    return WorkflowAgent(
        workflow=workflow,
        name="MedTechSalesWorkflow",
        description="End-to-end: transcript analysis → presentation → email",
    )


# ── Follow-Up Workflow (generate draft, don't send) ──
def build_followup_agent() -> Agent:
    """Presentation agent that drafts the email instead of sending it."""
    return Agent(
        client=get_foundry_client(),
        name="FollowUpPresentationAgent",
        instructions=(
            "You are a MedTech follow-up specialist.\n"
            "Based on the transcript analysis you receive:\n"
            "1. Create a professional follow-up presentation using create_presentation.\n"
            "2. Prepare an email draft using prepare_email_draft — do NOT send the email.\n"
            "   The email should be a professional follow-up to the surgeon referencing the call highlights.\n"
            "   If no recipient email is specified, use 'recipient@example.com' as placeholder.\n\n"
            "Always create the presentation first, then prepare the email draft."
        ),
        tools=[create_presentation, prepare_email_draft],
    )


def build_followup_workflow() -> WorkflowAgent:
    transcript_agent = build_transcript_agent()
    followup_agent = build_followup_agent()

    workflow = (
        WorkflowBuilder(
            name="TranscriptToFollowUp",
            description="Analyzes a MedTech sales call transcript, creates a presentation and prepares an email draft for review.",
            start_executor=transcript_agent,
        )
        .add_edge(transcript_agent, followup_agent)
        .build()
    )

    return WorkflowAgent(
        workflow=workflow,
        name="MedTechFollowUpWorkflow",
        description="End-to-end: transcript analysis → presentation → email draft (no send)",
    )


# ── Podcast: SSML script writer ──
PODCAST_VOICES = {
    "host": {"name": "Alex", "role": "Host", "voice": "en-US-AndrewMultilingualNeural"},
    "clinical": {"name": "Dr. Sarah", "role": "Clinical Expert", "voice": "en-US-AvaMultilingualNeural"},
    "sales": {"name": "Mike", "role": "Sales Specialist", "voice": "en-US-BrianMultilingualNeural"},
}


def build_podcast_agent(product_key: str = None) -> Agent:
    product = PRODUCTS.get(product_key) if product_key else None

    product_context = ""
    if product:
        product_context = (
            f"\n\n--- PRODUCT TO DISCUSS ---\n"
            f"Product: {product['name']}\n"
            f"{product['description']}\n"
        )

    host = PODCAST_VOICES["host"]
    clinical = PODCAST_VOICES["clinical"]
    sales = PODCAST_VOICES["sales"]

    instructions = (
        "You are a podcast script writer for MedTech Insights, a professional medical technology podcast.\n\n"
        "Your task is to write a podcast episode script as valid SSML (Speech Synthesis Markup Language) "
        "that can be sent directly to Azure Speech SDK for audio synthesis.\n\n"
        "The podcast has 3 characters:\n"
        f"1. **{host['name']}** ({host['role']}) — The engaging host who introduces the product, "
        f"asks insightful questions, and guides the conversation. Voice: {host['voice']}\n"
        f"2. **{clinical['name']}** ({clinical['role']}) — A knowledgeable surgeon/clinician who "
        f"discusses clinical benefits, evidence, and patient outcomes. Voice: {clinical['voice']}\n"
        f"3. **{sales['name']}** ({sales['role']}) — An experienced MedTech sales professional who "
        f"covers practical adoption, OR workflow, value proposition, and competitive positioning. Voice: {sales['voice']}\n\n"
        "SSML FORMAT RULES — follow these EXACTLY:\n"
        "- Output ONLY the SSML markup. No extra text, no markdown code blocks.\n"
        "- Start with: <speak version=\"1.0\" xmlns=\"http://www.w3.org/2001/10/synthesis\" xml:lang=\"en-US\">\n"
        "- End with: </speak>\n"
        "- Each character's dialogue is wrapped in a <voice name=\"...\"> element.\n"
        "- CRITICAL: ALL elements (<break>, <prosody>, <emphasis>) MUST be INSIDE a <voice> element. "
        "Nothing except <voice> elements may appear as direct children of <speak>.\n"
        "- To add a pause between speakers, put <break time=\"500ms\"/> at the END of the previous speaker's <voice> block "
        "(before the closing </voice> tag). NEVER place <break> between </voice> and <voice>.\n"
        "- Use <prosody> and <emphasis> inside <voice> elements only.\n"
        "- Use <emphasis level=\"moderate\"> sparingly to highlight key product names or features.\n"
        "- Use <break time=\"300ms\"/> for brief pauses within a speaker's dialogue.\n"
        "- Do NOT use any characters that would break XML: use &amp; for &, &lt; for <, &gt; for >.\n"
        "- Do NOT include any text outside the <speak> element.\n\n"
        "SCRIPT STRUCTURE:\n"
        "1. Opening: Host introduces the podcast and today's topic (~2-3 sentences)\n"
        "2. Product Overview: Host asks clinical expert to explain the product (~3-4 exchanges)\n"
        "3. Clinical Benefits: Deep dive into clinical evidence and patient outcomes (~3-4 exchanges)\n"
        "4. Sales Perspective: How reps can position the product, handle objections (~3-4 exchanges)\n"
        "5. Closing: Host summarizes key takeaways, each guest gives one final tip (~2-3 sentences each)\n\n"
        "Keep the total script to about 15-20 speaker turns. Each turn should be 2-4 sentences. "
        "Maintain a professional but conversational tone — informative and engaging, like a real medical podcast."
        f"{product_context}"
    )

    return Agent(
        client=get_foundry_client(),
        name="PodcastScriptWriterAgent",
        instructions=instructions,
    )


async def main() -> None:
    workflow_agent = build_workflow()

    # <run_workflow>
    result = await workflow_agent.run(
        "Here is the call transcript:\n\n"
        "Sales Rep: Hi Dr. Patel, thanks for taking the time today. I wanted to walk you through our new surgical robot, the MedAssist Pro.\n"
        "Surgeon: Sure. I've been hearing about it. What's the learning curve?\n"
        "Sales Rep: Most surgeons get comfortable within 5-7 cases. We also offer on-site training.\n"
        "Surgeon: That's reasonable. What about integration with our existing OR setup?\n"
        "Sales Rep: It's compatible with all major imaging systems. We can do a site assessment.\n"
        "Surgeon: I'd like that. Also, the cost — can you share a breakdown?\n"
        "Sales Rep: Absolutely, I'll send over a detailed pricing proposal by Friday.\n"
        "Surgeon: Great. Let me also discuss this with our department head, Dr. Kim.\n"
        "Sales Rep: Perfect. Should I set up a demo for the both of you?\n"
        "Surgeon: Yes, sometime next week would work.\n\n"
        "Email the presentation to mukulsaluja@microsoft.com."
    )
    print(f"Workflow result: {result}")
    # </run_workflow>


if __name__ == "__main__":
    asyncio.run(main())